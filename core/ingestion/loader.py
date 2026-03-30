"""
Document loader — reads PDF, DOCX, PPTX, and TXT files into LangChain Document objects.
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import List

from langchain_core.documents import Document


def load_pdf(file_path: str) -> List[Document]:
    """Load a PDF file and return a list of Documents (one per page)."""
    import PyPDF2

    documents: List[Document] = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                documents.append(
                    Document(
                        page_content=text,
                        metadata={
                            "source": os.path.basename(file_path),
                            "page": i + 1,
                            "file_path": file_path,
                        },
                    )
                )
    return documents


def load_docx(file_path: str) -> List[Document]:
    """Load a DOCX file and return a list of Documents."""
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    full_text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    if not full_text.strip():
        return []
    return [
        Document(
            page_content=full_text,
            metadata={
                "source": os.path.basename(file_path),
                "file_path": file_path,
            },
        )
    ]


def load_pptx(file_path: str) -> List[Document]:
    """Load a PPTX file and return a list of Documents (one per slide)."""
    from pptx import Presentation

    prs = Presentation(file_path)
    documents: List[Document] = []
    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slide_text = "\n".join(texts)
        if slide_text.strip():
            documents.append(
                Document(
                    page_content=slide_text,
                    metadata={
                        "source": os.path.basename(file_path),
                        "slide": i + 1,
                        "file_path": file_path,
                    },
                )
            )
    return documents


def load_txt(file_path: str) -> List[Document]:
    """Load a plain-text file and return a single Document."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    if not text.strip():
        return []
    return [
        Document(
            page_content=text,
            metadata={
                "source": os.path.basename(file_path),
                "file_path": file_path,
            },
        )
    ]


# Mapping of extension → loader function
_LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".txt": load_txt,
}


def load_document(file_path: str) -> List[Document]:
    """
    Auto-detect file type and load into Documents.

    Raises:
        ValueError: If the file extension is not supported.
    """
    ext = Path(file_path).suffix.lower()
    loader_fn = _LOADERS.get(ext)
    if loader_fn is None:
        raise ValueError(
            f"Unsupported file type: {ext}. Supported: {list(_LOADERS.keys())}"
        )
    return loader_fn(file_path)
